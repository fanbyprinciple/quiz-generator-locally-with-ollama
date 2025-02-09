import streamlit as st
import json
import ollama
from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import io
import re

@st.cache_data
def fetch_ppt_content(text_content, slide_count):
    RESPONSE_JSON = {
        "slides": [
            {
                "title": "Clear Slide Title 1",
                "content": ["Concise point 1", "Relevant point 2", "Key takeaway 3"]
            }
        ]
    }

    PROMPT_TEMPLATE = f"""
    Text: {text_content}
    Create a structured PowerPoint presentation with exactly {slide_count} slides.
    Follow these rules:
    1. Slide titles should be 3-7 words
    2. Each slide should have 3-5 bullet points
    3. Content must be extracted from the text
    4. Use professional business language
    
    Format your response EXACTLY like this JSON:
    {json.dumps(RESPONSE_JSON, indent=2)}
    Replace the example content with real content from the text.
    """

    response = ollama.chat(
        model='llama3.2:3b',
        messages=[{
            'role': 'user',
            'content': PROMPT_TEMPLATE
        }],
        format='json',
        options={
            'temperature': 0.2,
            'num_ctx': 2048
        }
    )

    try:
        content = json.loads(response['message']['content'])
        if len(content["slides"]) != slide_count:
            st.error(f"Requested {slide_count} slides but got {len(content['slides'])}")
            return None
        return content["slides"]
    except Exception as e:
        st.error(f"Failed to parse response: {str(e)}")
        return None

def extract_text_from_file(uploaded_file):
    """Extract text from uploaded PDF or text files."""
    if uploaded_file.name.endswith('.pdf'):
        pdf_reader = PdfReader(uploaded_file)
        return " ".join(page.extract_text() or "" for page in pdf_reader.pages).strip()
    elif uploaded_file.name.endswith(('.txt', '.md')):
        return uploaded_file.getvalue().decode("utf-8").strip()
    else:
        st.error("Unsupported file format")
        return None

def create_ppt_from_template(slides_data, template_file):
    """Generate a PPT using the provided template for design only."""
    prs = Presentation(template_file)
    
    # Extract background color if available
    try:
        template_slide = prs.slides[0]
        background_fill = template_slide.background.fill
        background_fill.solid()  # Ensure solid fill before setting color
        background_color = background_fill.fore_color.rgb
    except AttributeError:
        background_color = None  # If no valid background, skip setting it

    # Create a blank presentation for content
    prs_new = Presentation()
    
    # Apply background color if available
    for slide_master in prs_new.slide_masters:
        slide_master.background.fill.solid()
        if background_color:
            slide_master.background.fill.fore_color.rgb = background_color

    # Add slides with content
    for slide_data in slides_data:
        slide_layout = prs_new.slide_layouts[5]  # Use Blank layout
        slide = prs_new.slides.add_slide(slide_layout)

        # Dynamically choose text color based on background
        text_color = RGBColor(255, 255, 255) if background_color and sum(background_color) < 382 else RGBColor(0, 0, 0)

        # Add title (ensure it exists)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = sanitize_text(slide_data["title"])
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.color.rgb = text_color

        # Add content
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame

        for point in slide_data["content"]:
            p = text_frame.add_paragraph()
            p.text = sanitize_text(point)
            p.level = 0
            p.font.size = Pt(20)
            p.font.color.rgb = text_color  # Ensure text visibility

    # Save to bytes buffer for Streamlit
    ppt_buffer = io.BytesIO()
    prs_new.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer

def sanitize_text(text):
    """Clean up model-generated text."""
    text = re.sub(r'[\x00-\x1F\x7F-\x9F]', '', text)  # Remove control characters
    text = text.strip('•-* ')  # Remove leading/trailing bullets
    return text[:500]  # Limit text length

def main():
    st.title("Professional PPT Generator")

    # File upload section
    col1, col2 = st.columns(2)
    with col1:
        uploaded_file = st.file_uploader("Upload document (PDF/TXT)", type=['pdf', 'txt', 'md'])
    with col2:
        template_file = st.file_uploader("Upload PPT template (design only)", type=['pptx'])

    # Configuration
    slide_count = st.number_input("Number of slides to generate", min_value=1, max_value=20, value=5)

    if st.button("Generate Presentation") and uploaded_file and template_file:
        with st.spinner("Analyzing document and generating slides..."):
            # Extract text
            text_content = extract_text_from_file(uploaded_file)
            
            if text_content:
                # Generate slide content
                slides_data = fetch_ppt_content(text_content, slide_count)
                
                if slides_data:
                    # Create PPT from template
                    try:
                        ppt_buffer = create_ppt_from_template(slides_data, template_file)
                        if ppt_buffer:
                            st.success("Presentation created successfully!")
                            
                            # Download button
                            st.download_button(
                                label="Download PowerPoint",
                                data=ppt_buffer,
                                file_name="generated_presentation.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                    except Exception as e:
                        st.error(f"PPT creation failed: {str(e)}")

if __name__ == "__main__":
    main()
